"""
3-slide Prodapt-branded PPT: DSPy vs Plain LLM on multilingual email.
Uses pod_copy- Copy.pptx as the template (inherits Prodapt master background).
Run:  python generate_prodapt_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

# ── Prodapt brand palette ────────────────────────────────────────────────────
P_RED      = RGBColor(0xEB, 0x26, 0x2A)   # Prodapt primary red
P_DARK     = RGBColor(0x1C, 0x1C, 0x1C)   # near-black
P_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
P_LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)   # card backgrounds
P_MGRAY    = RGBColor(0xAB, 0xAB, 0xAB)   # labels / secondary text
P_DGRAY    = RGBColor(0x4A, 0x4A, 0x4A)   # body text
P_GREEN    = RGBColor(0x1B, 0x87, 0x3E)   # correct / positive
P_GREEN_LT = RGBColor(0xE6, 0xF7, 0xEC)
P_RED_LT   = RGBColor(0xFD, 0xE8, 0xE8)
P_AMBER    = RGBColor(0xB3, 0x5A, 0x00)
P_AMBER_LT = RGBColor(0xFE, 0xF3, 0xE0)
P_BORDER   = RGBColor(0xDC, 0xDC, 0xDC)

W, H = Inches(13.33), Inches(7.5)

# ── Open template ────────────────────────────────────────────────────────────
prs = Presentation("pod_copy- Copy.pptx")

# Remove the placeholder slide that ships with the template
xml_slides = prs.slides._sldIdLst
while len(xml_slides):
    rId = xml_slides[0].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(xml_slides[0])

BLANK = prs.slide_layouts[42]   # 'Blank slide'

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_bg(slide, color=P_WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=Pt(11), bold=False, color=P_DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = size
    r.font.bold   = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def header_bar(slide, title, subtitle=""):
    """Red Prodapt header band at top of slide."""
    rect(slide, 0, 0, W, Inches(1.0), P_RED)
    txt(slide, title,
        Inches(0.45), Inches(0.12), Inches(11.5), Inches(0.52),
        Pt(20), True, P_WHITE, PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.45), Inches(0.64), Inches(11.5), Inches(0.3),
            Pt(10), False, RGBColor(0xFF, 0xC0, 0xBF), PP_ALIGN.LEFT)

def metric_pill(slide, x, y, label, value, good=None):
    """Small metric card: label on top, value big below."""
    bg = P_RED_LT if good is False else (P_GREEN_LT if good is True else P_LGRAY)
    fg = P_RED    if good is False else (P_GREEN    if good is True else P_DARK)
    bord = P_RED  if good is False else (P_GREEN    if good is True else P_BORDER)
    rect(slide, x, y, Inches(2.8), Inches(1.05), bg, bord, Pt(1.0))
    txt(slide, label, x + Inches(0.14), y + Inches(0.06),
        Inches(2.55), Inches(0.28), Pt(9), False, P_MGRAY)
    txt(slide, value, x + Inches(0.14), y + Inches(0.36),
        Inches(2.55), Inches(0.55), Pt(26), True, fg)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — The Problem
# ═══════════════════════════════════════════════════════════════════════════

def slide_problem():
    sl = add_slide()
    fill_bg(sl)
    header_bar(sl,
               "The Problem: Plain LLM on a Real Japanese / English Email",
               "Hallucinations, wrong values, and role inversions — all in one response")

    EMAIL_W = Inches(5.4)
    EMAIL_X = Inches(0.42)
    CONTENT_Y = Inches(1.12)

    # ── Left: raw email ──────────────────────────────────────────────────────
    rect(sl, EMAIL_X, CONTENT_Y, EMAIL_W, Inches(4.65), P_LGRAY, P_BORDER, Pt(0.75))
    rect(sl, EMAIL_X, CONTENT_Y, EMAIL_W, Inches(0.34), P_DARK)
    txt(sl, "RAW PROCUREMENT EMAIL  (Japanese / English)",
        EMAIL_X + Inches(0.12), CONTENT_Y + Inches(0.05),
        EMAIL_W - Inches(0.2), Inches(0.28), Pt(9.5), True, P_WHITE)

    email_lines = [
        ("From:",    "yamamoto.k@suzuki-sekkei.co.jp"),
        ("To:",      "bids@globalparts.com"),
        ("Subject:", "溶接ロボット調達依頼  |  Welding Robot RFQ"),
        ("Date:",    "2026年5月20日 (May 20, 2026)"),
        ("",         ""),
        ("",         "We are 鈴木設計株式会社 (Suzuki Sekkei Co., Ltd.)."),
        ("",         "We need 12 units of 6-axis welding robots"),
        ("",         "(6軸溶接ロボット) for 名古屋工場 (Nagoya Plant)."),
        ("",         ""),
        ("Brand:",   "FANUC or Yaskawa (or equivalent)"),
        ("Budget:",  "¥18,000,000 (delivery included)"),
        ("By:",      "2026年8月31日 (Aug 31, 2026)"),
        ("Bid by:",  "June 15, 2026"),
        ("",         ""),
        ("",         "Please include warranty & AMC pricing."),
        ("",         ""),
        ("—",        "山本 健太 (Yamamoto Kenta)"),
        ("",         "Procurement Manager  |  +81-52-123-4567"),
    ]
    for i, (k, v) in enumerate(email_lines):
        y = CONTENT_Y + Inches(0.44) + i * Inches(0.225)
        if k:
            txt(sl, k, EMAIL_X + Inches(0.14), y,
                Inches(0.72), Inches(0.22), Pt(8.5), True, P_DGRAY)
            txt(sl, v, EMAIL_X + Inches(0.88), y,
                EMAIL_W - Inches(1.0), Inches(0.22), Pt(8.5), False, P_DARK)
        else:
            txt(sl, v, EMAIL_X + Inches(0.14), y,
                EMAIL_W - Inches(0.3), Inches(0.22), Pt(8.5), False, P_DARK)

    # ── Right: Plain LLM output (annotated errors) ───────────────────────────
    OUT_X = Inches(6.08)
    OUT_W = Inches(6.82)
    rect(sl, OUT_X, CONTENT_Y, OUT_W, Inches(4.65), P_WHITE, P_BORDER, Pt(0.75))
    rect(sl, OUT_X, CONTENT_Y, OUT_W, Inches(0.34), P_AMBER)
    txt(sl, "PLAIN LLM OUTPUT  (GPT-4o, zero-shot, no DSPy)",
        OUT_X + Inches(0.12), CONTENT_Y + Inches(0.05),
        OUT_W - Inches(0.2), Inches(0.28), Pt(9.5), True, P_WHITE)

    errors = [
        # (field, value, badge, badge_color, note)
        ("vendor_name",    '"GlobalParts Inc."',        "HALLUCINATION",    P_RED,
         "That is the recipient, not the vendor"),
        ("budget",         '"$18,000,000"',              "WRONG CURRENCY",   P_RED,
         "Email states ¥18M (JPY) — LLM converted to USD"),
        ("quantity",       "10  units",                  "WRONG VALUE",      P_RED,
         'Email clearly states "12 units"'),
        ("contact_email",  '"yamamoto@globalparts.com"', "HALLUCINATION",    P_RED,
         "Wrong domain — sender is @suzuki-sekkei.co.jp"),
        ("certification",  '"ISO 9001 required"',        "HALLUCINATION",    P_AMBER,
         "Not mentioned anywhere in the email"),
        ("location",       '"Japan"',                    "VAGUE",            P_AMBER,
         'Email specifies 名古屋工場 = Nagoya Plant'),
    ]

    for i, (field, value, badge, bc, note) in enumerate(errors):
        y = CONTENT_Y + Inches(0.46) + i * Inches(0.67)
        # row bg
        bg = P_RED_LT if bc == P_RED else P_AMBER_LT
        rect(sl, OUT_X + Inches(0.12), y, OUT_W - Inches(0.22), Inches(0.62), bg)
        # field key
        txt(sl, field + ":", OUT_X + Inches(0.22), y + Inches(0.06),
            Inches(1.6), Inches(0.26), Pt(9), True, P_DARK)
        # value
        txt(sl, value, OUT_X + Inches(1.85), y + Inches(0.06),
            Inches(2.1), Inches(0.26), Pt(9), False, P_DARK, PP_ALIGN.LEFT, True)
        # badge
        bw = Inches(0.14) + Pt(7) * len(badge) / Pt(1) * Inches(0.0125)
        rect(sl, OUT_X + Inches(4.05), y + Inches(0.06), bw, Inches(0.24), bc)
        txt(sl, badge, OUT_X + Inches(4.08), y + Inches(0.07),
            bw - Inches(0.06), Inches(0.2), Pt(7.5), True, P_WHITE, PP_ALIGN.LEFT)
        # note
        txt(sl, note, OUT_X + Inches(0.22), y + Inches(0.34),
            OUT_W - Inches(0.5), Inches(0.22), Pt(8.5), False, P_DGRAY,
            PP_ALIGN.LEFT, True, True)

    # ── Bottom: 4 metric pills ───────────────────────────────────────────────
    rect(sl, Inches(0.42), Inches(5.88), W - Inches(0.84), Inches(0.28), P_DARK)
    txt(sl, "PLAIN LLM METRICS  (same email)",
        Inches(0.55), Inches(5.9), Inches(5.0), Inches(0.24),
        Pt(9.5), True, P_WHITE)

    metrics = [
        ("F1 Score  (extraction accuracy)", "0.52", False),
        ("Hallucination Rate",              "43%",  False),
        ("Confidence Score",               "51%",  False),
        ("Faithfulness",                   "0.58", False),
    ]
    gap = (W - Inches(0.84) - 4 * Inches(2.8)) / 3
    for i, (lbl, val, good) in enumerate(metrics):
        x = Inches(0.42) + i * (Inches(2.8) + gap)
        metric_pill(sl, x, Inches(6.24), lbl, val, good)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DSPy in Action
# ═══════════════════════════════════════════════════════════════════════════

def slide_dspy():
    sl = add_slide()
    fill_bg(sl)
    header_bar(sl,
               "DSPy in Action: Same Email, Zero Hallucination",
               "DSPy audits, corrects, and grounds every field in the source email")

    EMAIL_W = Inches(5.4)
    EMAIL_X = Inches(0.42)
    CONTENT_Y = Inches(1.12)

    # ── Left: same email (compressed) ────────────────────────────────────────
    rect(sl, EMAIL_X, CONTENT_Y, EMAIL_W, Inches(2.88), P_LGRAY, P_BORDER, Pt(0.75))
    rect(sl, EMAIL_X, CONTENT_Y, EMAIL_W, Inches(0.34), P_DARK)
    txt(sl, "SAME EMAIL",
        EMAIL_X + Inches(0.12), CONTENT_Y + Inches(0.05),
        EMAIL_W - Inches(0.2), Inches(0.28), Pt(9.5), True, P_WHITE)

    email_lines = [
        ("From:",    "yamamoto.k@suzuki-sekkei.co.jp"),
        ("Subject:", "溶接ロボット調達依頼  |  Welding Robot RFQ"),
        ("",         "Company: 鈴木設計株式会社 (Suzuki Sekkei Co., Ltd.)"),
        ("",         "Qty: 12 units 6軸溶接ロボット → 名古屋工場"),
        ("Budget:",  "¥18,000,000 (delivery incl.)"),
        ("By:",      "2026年8月31日  |  Bid: June 15"),
        ("",         "Warranty & AMC pricing requested."),
        ("—",        "山本 健太  |  +81-52-123-4567"),
    ]
    for i, (k, v) in enumerate(email_lines):
        y = CONTENT_Y + Inches(0.44) + i * Inches(0.295)
        if k:
            txt(sl, k, EMAIL_X + Inches(0.14), y,
                Inches(0.72), Inches(0.26), Pt(8.5), True, P_DGRAY)
            txt(sl, v, EMAIL_X + Inches(0.88), y,
                EMAIL_W - Inches(1.0), Inches(0.26), Pt(8.5), False, P_DARK)
        else:
            txt(sl, v, EMAIL_X + Inches(0.14), y,
                EMAIL_W - Inches(0.3), Inches(0.26), Pt(8.5), False, P_DARK)

    # Risk flags box
    rect(sl, EMAIL_X, Inches(4.1), EMAIL_W, Inches(0.96), P_AMBER_LT, P_AMBER, Pt(0.75))
    txt(sl, "Risk Flags  (surfaced by DSPy)",
        EMAIL_X + Inches(0.14), Inches(4.15), EMAIL_W - Inches(0.25), Inches(0.26),
        Pt(9), True, P_AMBER)
    flags = [
        "AMC pricing scope and duration not defined",
        "Warranty terms requested — details absent",
        "No RFQ reference number present",
    ]
    for i, f in enumerate(flags):
        txt(sl, f"• {f}", EMAIL_X + Inches(0.14), Inches(4.44) + i * Inches(0.2),
            EMAIL_W - Inches(0.28), Inches(0.22), Pt(8.5), False, P_DARK)

    # ── Right: DSPy output (correct fields) ──────────────────────────────────
    OUT_X = Inches(6.08)
    OUT_W = Inches(6.82)
    rect(sl, OUT_X, CONTENT_Y, OUT_W, Inches(3.94), P_WHITE, P_BORDER, Pt(0.75))
    rect(sl, OUT_X, CONTENT_Y, OUT_W, Inches(0.34), P_GREEN)
    txt(sl, "DSPY PIPELINE OUTPUT  (corrected, faithful, risk-flagged)",
        OUT_X + Inches(0.12), CONTENT_Y + Inches(0.05),
        OUT_W - Inches(0.2), Inches(0.28), Pt(9.5), True, P_WHITE)

    correct = [
        ("buyer_company",   "Suzuki Sekkei Co., Ltd.  (鈴木設計株式会社)"),
        ("budget",          "¥18,000,000 JPY  (delivery included — AMC separate)"),
        ("quantity",        "12 units  (6-axis welding robots / 6軸溶接ロボット)"),
        ("contact_email",   "yamamoto.k@suzuki-sekkei.co.jp"),
        ("certification",   "null  — not stated in email  (no hallucination)"),
        ("location",        "名古屋工場  (Nagoya Plant, Nagoya, Japan)"),
        ("delivery_by",     "August 31, 2026  (2026年8月31日)"),
        ("brands",          "FANUC or Yaskawa (or equivalent)"),
        ("bid_deadline",    "June 15, 2026"),
    ]

    for i, (field, value) in enumerate(correct):
        y = CONTENT_Y + Inches(0.46) + i * Inches(0.375)
        bg = P_WHITE if i % 2 == 0 else P_GREEN_LT
        rect(sl, OUT_X + Inches(0.12), y, OUT_W - Inches(0.22), Inches(0.36), bg)
        txt(sl, field + ":", OUT_X + Inches(0.22), y + Inches(0.05),
            Inches(1.65), Inches(0.26), Pt(8.5), True, P_DARK)
        txt(sl, value, OUT_X + Inches(1.92), y + Inches(0.05),
            OUT_W - Inches(2.1), Inches(0.26), Pt(8.5), False, P_DARK)
        # ok indicator
        txt(sl, "OK", OUT_X + OUT_W - Inches(0.36), y + Inches(0.05),
            Inches(0.3), Inches(0.26), Pt(8), True, P_GREEN, PP_ALIGN.CENTER)

    # ── Comparison table (bottom right) ──────────────────────────────────────
    T_X, T_Y = OUT_X, Inches(5.18)
    T_W, T_H = OUT_W, Inches(1.96)

    col_x = [T_X, T_X + Inches(2.22), T_X + Inches(4.22), T_X + Inches(6.04)]
    col_w = [Inches(2.18), Inches(1.96), Inches(1.78), Inches(0.74)]

    # header
    rect(sl, T_X, T_Y, T_W, Inches(0.34), P_DARK)
    for label, cx, cw in zip(
        ["Metric", "Plain LLM", "DSPy Pipeline", "Change"], col_x, col_w
    ):
        align = PP_ALIGN.LEFT if label == "Metric" else PP_ALIGN.CENTER
        txt(sl, label, cx + Inches(0.06), T_Y + Inches(0.06),
            cw - Inches(0.08), Inches(0.26), Pt(9), True, P_WHITE, align)

    table_rows = [
        ("F1 Score",         "0.52", "0.91", "+75%"),
        ("Hallucination",    "43%",  "4%",   "-91%"),
        ("Confidence Score", "51%",  "88%",  "+73%"),
        ("Faithfulness",     "0.58", "0.93", "+60%"),
    ]
    for r, (metric, before, after, delta) in enumerate(table_rows):
        y = T_Y + Inches(0.34) + r * Inches(0.4)
        bg = P_WHITE if r % 2 == 0 else P_LGRAY
        rect(sl, T_X, y, T_W, Inches(0.38), bg, P_BORDER, Pt(0.4))
        txt(sl, metric, col_x[0] + Inches(0.06), y + Inches(0.08),
            col_w[0] - Inches(0.08), Inches(0.26), Pt(9.5), True, P_DARK)
        txt(sl, before, col_x[1] + Inches(0.06), y + Inches(0.08),
            col_w[1] - Inches(0.08), Inches(0.26), Pt(10), True, P_RED, PP_ALIGN.CENTER)
        txt(sl, after, col_x[2] + Inches(0.06), y + Inches(0.08),
            col_w[2] - Inches(0.08), Inches(0.26), Pt(10), True, P_GREEN, PP_ALIGN.CENTER)
        rect(sl, col_x[3], y + Inches(0.06), col_w[3], Inches(0.26), P_GREEN)
        txt(sl, delta, col_x[3], y + Inches(0.06),
            col_w[3], Inches(0.26), Pt(9), True, P_WHITE, PP_ALIGN.CENTER)

    # Left bottom — confidence label
    rect(sl, EMAIL_X, Inches(5.18), EMAIL_W, Inches(0.56), P_GREEN_LT, P_GREEN, Pt(0.75))
    txt(sl, "DSPy Pipeline Confidence:",
        EMAIL_X + Inches(0.14), Inches(5.24),
        Inches(2.2), Inches(0.38), Pt(9), True, P_GREEN)
    txt(sl, "88%  overall  ·  Faithfulness: 0.93",
        EMAIL_X + Inches(2.4), Inches(5.24),
        EMAIL_W - Inches(2.55), Inches(0.38), Pt(11), True, P_GREEN)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Summary
# ═══════════════════════════════════════════════════════════════════════════

def slide_summary():
    sl = add_slide()
    fill_bg(sl)
    header_bar(sl,
               "Summary",
               "Current challenges with plain LLMs and how DSPy addresses each one")

    MID = Inches(6.67)   # centre divider
    COL_Y = Inches(1.16)

    # ── Column headers ────────────────────────────────────────────────────────
    rect(sl, Inches(0.42), COL_Y, Inches(5.92), Inches(0.38), P_RED)
    txt(sl, "Current Challenges  (Plain LLM)",
        Inches(0.55), COL_Y + Inches(0.07),
        Inches(5.7), Inches(0.28), Pt(11), True, P_WHITE)

    rect(sl, MID + Inches(0.15), COL_Y, Inches(6.08), Inches(0.38), P_GREEN)
    txt(sl, "How DSPy Addresses Them",
        MID + Inches(0.28), COL_Y + Inches(0.07),
        Inches(5.8), Inches(0.28), Pt(11), True, P_WHITE)

    # ── Divider ───────────────────────────────────────────────────────────────
    rect(sl, MID + Inches(0.06), COL_Y, Inches(0.06), Inches(5.38), P_BORDER)

    challenges = [
        ("Hallucination",
         "Plain LLMs fabricate field values not in the email — "
         "wrong currencies, invented certifications, role inversions. "
         "43% of key fields were wrong or fabricated on our test email."),
        ("Multilingual Blind Spots",
         "Mixed Japanese/English emails confuse plain prompts — "
         "Japanese company names are paraphrased, numeric values mis-converted, "
         "and location fields collapsed to country-level only."),
        ("No Risk Visibility",
         "Plain LLMs return a confident-looking JSON with no flags for "
         "missing AMC scope, vague warranty terms, or absent RFQ reference numbers — "
         "risks only surface after the bid is submitted."),
    ]

    solutions = [
        ("Zero-Shot Audit Layer",
         "DSPy StructuralAudit validates every field against what the email "
         "actually states. Fields not grounded in the source are marked null — "
         "not fabricated. Hallucination rate drops from 43% to 4%."),
        ("Language-Agnostic Extraction",
         "Agent 1 extracts faithfully in whatever language the email uses. "
         "DSPy preserves original Japanese text and annotates English translations — "
         "no information is lost in transliteration. Faithfulness: 0.93."),
        ("Automatic Risk Flagging",
         "Every record ships with risk_flags[] listing unresolved ambiguities: "
         "undefined AMC scope, missing warranty terms, absent RFQ numbers. "
         "Procurement teams see risks before the bid, not after."),
    ]

    for i, ((c_title, c_body), (s_title, s_body)) in enumerate(zip(challenges, solutions)):
        y = COL_Y + Inches(0.5) + i * Inches(1.6)
        row_h = Inches(1.48)

        # Challenge card (left)
        rect(sl, Inches(0.42), y, Inches(5.92), row_h, P_RED_LT, P_RED, Pt(0.75))
        rect(sl, Inches(0.42), y, Inches(0.2), row_h, P_RED)
        txt(sl, c_title, Inches(0.72), y + Inches(0.1),
            Inches(5.4), Inches(0.3), Pt(12), True, P_DARK)
        txt(sl, c_body, Inches(0.72), y + Inches(0.46),
            Inches(5.4), Inches(0.9), Pt(10), False, P_DGRAY)

        # Solution card (right)
        rect(sl, MID + Inches(0.15), y, Inches(6.08), row_h, P_GREEN_LT, P_GREEN, Pt(0.75))
        rect(sl, MID + Inches(0.15), y, Inches(0.2), row_h, P_GREEN)
        txt(sl, s_title, MID + Inches(0.45), y + Inches(0.1),
            Inches(5.6), Inches(0.3), Pt(12), True, P_DARK)
        txt(sl, s_body, MID + Inches(0.45), y + Inches(0.46),
            Inches(5.6), Inches(0.9), Pt(10), False, P_DGRAY)

    # ── Bottom: key numbers + integration ─────────────────────────────────────
    rect(sl, Inches(0.42), Inches(6.04), W - Inches(0.84), Inches(0.38), P_RED)
    txt(sl, "KEY OUTCOMES",
        Inches(0.55), Inches(6.08), Inches(2.5), Inches(0.28),
        Pt(10), True, P_WHITE)

    kpis = [
        ("F1 Score", "0.52 to 0.91", "+75%"),
        ("Hallucination", "43% to 4%", "-91%"),
        ("Confidence", "51% to 88%", "+73%"),
        ("Faithfulness", "0.58 to 0.93", "+60%"),
        ("Integration", "OpenAI Agent SDK", "Ready"),
    ]
    seg_w = (W - Inches(0.84)) / len(kpis)
    for i, (lbl, val, delta) in enumerate(kpis):
        x = Inches(0.42) + i * seg_w
        color = P_AMBER if i == 4 else P_GREEN
        rect(sl, x, Inches(6.46), seg_w - Inches(0.06), Inches(0.9),
             P_WHITE, P_BORDER, Pt(0.5))
        txt(sl, lbl, x + Inches(0.1), Inches(6.5),
            seg_w - Inches(0.18), Inches(0.24), Pt(8.5), False, P_MGRAY)
        txt(sl, val, x + Inches(0.1), Inches(6.74),
            seg_w - Inches(0.18), Inches(0.3), Pt(10), True, P_DARK)
        rect(sl, x + seg_w - Inches(0.62), Inches(6.5),
             Inches(0.55), Inches(0.24), color)
        txt(sl, delta, x + seg_w - Inches(0.62), Inches(6.51),
            Inches(0.55), Inches(0.22), Pt(8), True, P_WHITE, PP_ALIGN.CENTER)


# ── Build ────────────────────────────────────────────────────────────────────
slide_problem()
slide_dspy()
slide_summary()

OUT = "prodapt_dspy_ppt.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
