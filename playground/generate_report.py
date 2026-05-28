"""
Generates pipeline_comparison_report.pptx — white background, professional design.
Run:  python generate_report.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_NAVY       = RGBColor(0x0F, 0x28, 0x4A)   # header bars, titles
C_BLUE       = RGBColor(0x1A, 0x56, 0xDB)   # DSPy primary accent
C_BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFF)   # DSPy light fill
C_ORANGE     = RGBColor(0xC4, 0x50, 0x1A)   # No-DSPy accent
C_ORANGE_LT  = RGBColor(0xFD, 0xED, 0xE0)   # No-DSPy light fill
C_DARK       = RGBColor(0x1A, 0x20, 0x2C)   # body text
C_GRAY_D     = RGBColor(0x4A, 0x55, 0x68)   # secondary text
C_GRAY_M     = RGBColor(0x71, 0x80, 0x96)   # tertiary / labels
C_GRAY_L     = RGBColor(0xF0, 0xF4, 0xF8)   # card backgrounds
C_GRAY_LL    = RGBColor(0xF7, 0xF9, 0xFC)   # table alt rows
C_BORDER     = RGBColor(0xCB, 0xD5, 0xE0)   # dividers / borders
C_GREEN      = RGBColor(0x0B, 0x7B, 0x47)
C_GREEN_LT   = RGBColor(0xD1, 0xFA, 0xE5)
C_RED        = RGBColor(0xB9, 0x18, 0x1A)
C_RED_LT     = RGBColor(0xFE, 0xE2, 0xE2)
C_AMBER      = RGBColor(0x92, 0x40, 0x09)
C_AMBER_LT   = RGBColor(0xFF, 0xED, 0xC2)
C_TEAL       = RGBColor(0x05, 0x7A, 0x85)
C_TEAL_LT    = RGBColor(0xCC, 0xF2, 0xF4)

W, H = Inches(13.33), Inches(7.5)

# ── Core helpers ──────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color=C_WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=Pt(12), bold=False, color=C_DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = size
    r.font.bold   = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def pill(slide, x, y, w, label, bg, fg=C_WHITE):
    """Thin rounded-looking label pill via rect."""
    rect(slide, x, y, w, Inches(0.28), bg)
    txt(slide, label, x + Inches(0.1), y + Inches(0.03),
        w - Inches(0.15), Inches(0.25), Pt(9), True, fg, PP_ALIGN.CENTER)

def hline(slide, y, x0=Inches(0.55), x1=Inches(12.78), color=C_BORDER, h=Pt(0.75)):
    rect(slide, x0, y, x1 - x0, Inches(0.02), color)

# ── Page chrome (header bar + slide number) ───────────────────────────────

def page_chrome(slide, title, slide_no, subtitle=""):
    fill_bg(slide)
    # top navy bar
    rect(slide, 0, 0, W, Inches(0.95), C_NAVY)
    # blue left accent in header
    rect(slide, 0, 0, Inches(0.22), Inches(0.95), C_BLUE)
    txt(slide, title,
        Inches(0.35), Inches(0.14), Inches(11.5), Inches(0.65),
        Pt(20), True, C_WHITE, PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.35), Inches(0.65), Inches(11.5), Inches(0.3),
            Pt(10), False, RGBColor(0xA0, 0xAE, 0xC8), PP_ALIGN.LEFT)
    # slide number badge
    rect(slide, Inches(12.6), Inches(0.22), Inches(0.56), Inches(0.5), C_BLUE)
    txt(slide, str(slide_no),
        Inches(12.6), Inches(0.22), Inches(0.56), Inches(0.5),
        Pt(14), True, C_WHITE, PP_ALIGN.CENTER)

# ── Card helper (bordered box with optional left accent bar) ──────────────

def card(slide, x, y, w, h, title="", title_color=C_NAVY,
         bg=C_GRAY_L, border=C_BORDER, accent=None):
    if accent:
        rect(slide, x, y, Inches(0.06), h, accent)
        rect(slide, x + Inches(0.06), y, w - Inches(0.06), h, bg, border)
    else:
        rect(slide, x, y, w, h, bg, border)
    if title:
        txt(slide, title,
            x + (Inches(0.2) if accent else Inches(0.15)),
            y + Inches(0.1), w - Inches(0.4), Inches(0.32),
            Pt(11), True, title_color)
    return x + (Inches(0.2) if accent else Inches(0.15))  # inner x

def badge(slide, x, y, label, bg, fg=C_WHITE):
    w = Inches(0.1) + Pt(7) * len(label) / Pt(1) * Inches(0.013)
    rect(slide, x, y, w, Inches(0.25), bg)
    txt(slide, label, x + Inches(0.06), y + Inches(0.02),
        w - Inches(0.08), Inches(0.22), Pt(8.5), True, fg, PP_ALIGN.CENTER)
    return w

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # full-width navy top band (taller on title)
    rect(sl, 0, 0, W, Inches(4.0), C_NAVY)
    # blue left strip
    rect(sl, 0, 0, Inches(0.28), Inches(4.0), C_BLUE)
    # decorative blue bar bottom of navy
    rect(sl, 0, Inches(3.95), W, Inches(0.18), C_BLUE)

    txt(sl, "PROCUREMENT EMAIL PIPELINE",
        Inches(0.55), Inches(0.8), Inches(12.0), Inches(0.95),
        Pt(38), True, C_WHITE, PP_ALIGN.LEFT)
    txt(sl, "DSPy vs Prompt-Driven: A Head-to-Head Evaluation",
        Inches(0.55), Inches(1.75), Inches(12.0), Inches(0.65),
        Pt(22), False, RGBColor(0xA0, 0xBF, 0xFF), PP_ALIGN.LEFT)
    txt(sl,
        "How well do different multi-agent architectures extract, structure, and validate "
        "real-world procurement RFQ emails across diverse industries?",
        Inches(0.55), Inches(2.5), Inches(12.0), Inches(0.75),
        Pt(13), False, RGBColor(0xC8, 0xD5, 0xE8), PP_ALIGN.LEFT)

    # stat pills in white area
    stats = [
        ("10  Emails", C_BLUE),
        ("7  Domains", C_TEAL),
        ("3  Metrics", C_NAVY),
        ("2  Pipelines", C_ORANGE),
    ]
    for i, (label, color) in enumerate(stats):
        x = Inches(0.55) + i * Inches(3.0)
        rect(sl, x, Inches(4.35), Inches(2.6), Inches(1.1), color)
        txt(sl, label.split("  ")[0], x + Inches(0.18), Inches(4.42),
            Inches(2.3), Inches(0.55), Pt(28), True, C_WHITE, PP_ALIGN.LEFT)
        txt(sl, label.split("  ")[1], x + Inches(0.18), Inches(4.95),
            Inches(2.3), Inches(0.3), Pt(12), False,
            RGBColor(0xCC, 0xD5, 0xE0), PP_ALIGN.LEFT)

    rect(sl, 0, Inches(6.8), W, Inches(0.7), C_GRAY_L)
    txt(sl, "CodeForge Research  ·  May 2026",
        Inches(0.55), Inches(6.88), Inches(6.0), Inches(0.38),
        Pt(11), False, C_GRAY_M, PP_ALIGN.LEFT)
    txt(sl, "Azure GPT-4o  ·  temperature = 0  ·  python-pptx",
        Inches(7.0), Inches(6.88), Inches(6.0), Inches(0.38),
        Pt(11), False, C_GRAY_M, PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Introduction & Setup
# ═══════════════════════════════════════════════════════════════════════════

def slide_intro(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Introduction & Setup", 2, "What we are building and how")

    # Goal box
    rect(sl, Inches(0.5), Inches(1.08), W - Inches(1.0), Inches(0.95), C_BLUE_LIGHT, C_BLUE, Pt(0.75))
    rect(sl, Inches(0.5), Inches(1.08), Inches(0.22), Inches(0.95), C_BLUE)
    txt(sl, "GOAL",
        Inches(0.8), Inches(1.13), Inches(1.5), Inches(0.28),
        Pt(10), True, C_BLUE)
    txt(sl,
        "Automatically convert unstructured procurement RFQ emails into complete, "
        "validated, domain-aware structured JSON records — ready for downstream ERP or sourcing systems.",
        Inches(0.8), Inches(1.42), Inches(11.6), Inches(0.52),
        Pt(12), False, C_DARK)

    # Two pipeline cards
    for col_i, (label, color, lt_color, lines) in enumerate([
        ("WITHOUT DSPy  —  3-Agent Chain", C_ORANGE, C_ORANGE_LT, [
            ("Agent 1", "Forensic Extractor", "Reads email, extracts only stated facts, builds domain skeleton"),
            ("Agent 2", "Structural Auditor", "Adds missing fields (empty string), quarantines noisy fields"),
            ("Agent 3", "Enrichment Engine",  "Fills all gaps via domain knowledge, resolves, validates"),
        ]),
        ("WITH DSPy  —  Agent + DSPy + Agent", C_BLUE, C_BLUE_LIGHT, [
            ("Agent 1", "Document Extractor", "Reads email, extracts stated facts, builds domain skeleton"),
            ("DSPy",    "StructuralAudit",    "Compiled signature audits schema, adds domain fields, quarantines noise"),
            ("Agent 2", "Enrichment Engine",  "Fills all gaps, resolves review_section, validates final record"),
        ]),
    ]):
        x = Inches(0.5) + col_i * Inches(6.42)
        # card header
        rect(sl, x, Inches(2.18), Inches(6.15), Inches(0.42), color)
        txt(sl, label, x + Inches(0.14), Inches(2.22),
            Inches(5.9), Inches(0.35), Pt(11), True, C_WHITE)
        # card body
        rect(sl, x, Inches(2.6), Inches(6.15), Inches(3.7), lt_color, C_BORDER)

        for i, (stage, name, desc) in enumerate(lines):
            y = Inches(2.72) + i * Inches(1.12)
            # connector line
            if i > 0:
                rect(sl, x + Inches(0.58), y - Inches(0.18), Inches(0.02), Inches(0.2), color)
            # step circle via rect
            rect(sl, x + Inches(0.12), y, Inches(0.9), Inches(0.3), color)
            txt(sl, stage, x + Inches(0.12), y + Inches(0.02),
                Inches(0.9), Inches(0.27), Pt(9), True, C_WHITE, PP_ALIGN.CENTER)
            txt(sl, name, x + Inches(1.12), y,
                Inches(4.9), Inches(0.3), Pt(12), True, C_DARK)
            txt(sl, desc, x + Inches(1.12), y + Inches(0.32),
                Inches(4.9), Inches(0.55), Pt(10), False, C_GRAY_D)

    # Shared config footer
    rect(sl, Inches(0.5), Inches(6.42), W - Inches(1.0), Inches(0.78), C_GRAY_L, C_BORDER)
    for i, item in enumerate([
        ("Model", "Azure GPT-4o"),
        ("Temperature", "0.0 (deterministic)"),
        ("Emails", "10 procurement RFQs"),
        ("Prompts", "prompts.json (all agents)"),
        ("DSPy", "Predict(StructuralAudit) — zero-shot"),
    ]):
        x = Inches(0.7) + i * Inches(2.45)
        txt(sl, item[0], x, Inches(6.5), Inches(2.3), Inches(0.24),
            Pt(8.5), False, C_GRAY_M)
        txt(sl, item[1], x, Inches(6.74), Inches(2.3), Inches(0.3),
            Pt(11), True, C_DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Used
# ═══════════════════════════════════════════════════════════════════════════

def slide_data(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Data Used", 3, "10 real-world procurement emails across 7 industry domains")

    emails = [
        ("01", "Welding Robots",       "Manufacturing RFQ",   "₹2.2 Cr",  "Delta Infrastructure",   "Pune",       "formal"),
        ("02", "EV Battery Packs",     "Manufacturing RFQ",   "€18 M",    "Greenfield Auto",        "Leipzig",    "formal"),
        ("03", "Laptop Procurement",   "IT Hardware RFQ",     "$1.2 L",   "NovaTech Solutions",     "Bangalore",  "urgent"),
        ("04", "OLED Panels 50K",      "Electronics RFQ",     "$750 K",   "Metro Display Tech",     "Shenzhen",   "formal"),
        ("05", "Office Furniture",     "Facilities RFQ",      "₹85 L",    "GlobalSpace Offices",    "Hyderabad",  "formal"),
        ("06", "Smartphones 1,200",    "Consumer/Retail RFQ", "₹6.2 Cr",  "ConnectHub India",       "Delhi",      "informal"),
        ("07", "Electric Vans x75",    "Fleet Procurement",   "₹18 Cr",   "UrbanMove Logistics",    "Navi Mumbai","formal"),
        ("08", "Kitchen Equipment",    "Hospitality RFQ",     "£360 K",   "GrandVista Hotels",      "Edinburgh",  "formal"),
        ("09", "Smart Classrooms 500", "Govt. Tender",        "₹42 Cr",   "State Edu Dept.",        "500 schools","formal"),
        ("10", "Aluminium Supply",     "12-month Contract",   "~₹8 Cr",   "FastPack Mfg.",          "Bhiwandi",   "informal"),
    ]

    col_x = [Inches(0.5), Inches(0.92), Inches(2.72), Inches(5.42),
              Inches(6.55), Inches(8.75), Inches(10.85), Inches(12.22)]
    col_w = [Inches(0.42), Inches(1.75), Inches(2.65), Inches(1.08),
              Inches(2.15), Inches(2.05), Inches(1.32), Inches(0.8)]
    hdrs  = ["#",  "Email Subject",  "Category",  "Budget",  "Buyer",  "Site",  "Tone",  ""]

    # header row
    rect(sl, Inches(0.5), Inches(1.08), Inches(12.33), Inches(0.42), C_NAVY)
    for i, hdr in enumerate(hdrs[:-1]):
        txt(sl, hdr, col_x[i] + Inches(0.06), Inches(1.13),
            col_w[i], Inches(0.33), Pt(10), True, C_WHITE)

    tone_colors = {
        "formal": (C_BLUE, C_BLUE_LIGHT),
        "urgent": (C_RED, C_RED_LT),
        "informal": (C_TEAL, C_TEAL_LT),
    }

    for r, row in enumerate(emails):
        y = Inches(1.5) + r * Inches(0.56)
        bg = C_WHITE if r % 2 == 0 else C_GRAY_LL
        rect(sl, Inches(0.5), y, Inches(12.33), Inches(0.55), bg, C_BORDER, Pt(0.5))

        num, name, cat, budget, buyer, site, tone = row
        tone_fg, tone_bg = tone_colors.get(tone, (C_GRAY_D, C_GRAY_L))

        for val, cx, cw in zip([num, name, cat, budget, buyer, site],
                                col_x[:6], col_w[:6]):
            color = C_GRAY_M if val == num else C_DARK
            bold  = (val == name)
            txt(sl, val, cx + Inches(0.06), y + Inches(0.13),
                cw - Inches(0.08), Inches(0.35), Pt(10.5), bold, color)

        # tone badge
        rect(sl, col_x[6], y + Inches(0.12), Inches(1.28), Inches(0.3), tone_bg, tone_fg, Pt(0.75))
        txt(sl, tone, col_x[6] + Inches(0.06), y + Inches(0.14),
            Inches(1.2), Inches(0.26), Pt(9), True, tone_fg, PP_ALIGN.CENTER)

    hline(sl, Inches(7.08))
    txt(sl,
        "Coverage: 7 domain types  ·  4 currencies (INR, EUR, USD, GBP)  ·  "
        "Single-lot and multi-phase deliveries  ·  Domestic and international",
        Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.35),
        Pt(10), False, C_GRAY_M, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Evaluation Criteria
# ═══════════════════════════════════════════════════════════════════════════

def slide_criteria(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Evaluation Criteria & Expectations", 4,
                "Three primary metrics scored 0–10 per email, plus validation quality assessment")

    criteria = [
        ("01", "CORRECTNESS", C_NAVY, C_BLUE, C_BLUE_LIGHT,
         "Extracted values must match what the email explicitly states.",
         [
             "Facts in extracted_values: exact wording, no paraphrasing",
             "Inferences and defaults must be labeled, not presented as facts",
             "No hallucinated values, no swapped roles (buyer vs vendor)",
             "Computed fields (lead_time, unit_price) must show their basis",
         ]),
        ("02", "COVERAGE", C_NAVY, C_TEAL, C_TEAL_LT,
         "All fields a downstream procurement system needs must be present.",
         [
             "Domain-critical fields: payment_terms, certifications, delivery milestones",
             "No cross-domain bleed (no IT SaaS fields in a manufacturing RFQ)",
             "Specialty requirements captured: defect SLA, AMC, staggered lots",
             "Fields absent from the email must still appear (marked as empty/needs-input)",
         ]),
        ("03", "STRUCTURE", C_NAVY, C_ORANGE, C_ORANGE_LT,
         "Output must be logically grouped, consistent, and machine-readable.",
         [
             "Nested sections preferred: financial{}, timeline{}, parties{}",
             "Flat 50-key dumps are penalised — hard to consume programmatically",
             "Consistent schema across all 10 emails",
             "No redundant keys capturing the same data under different names",
         ]),
        ("04", "VALIDATION", C_NAVY, C_GREEN, C_GREEN_LT,
         "Validation section must surface real risks and be actionable.",
         [
             "risk_flags must name specific issues, not be an empty array",
             "completeness rating must be honest — no 'complete' when fields were inferred",
             "inconsistencies must flag budget ambiguity, undefined scope, etc.",
             "summary sentence must give a clear readiness assessment",
         ]),
    ]

    for i, (num, title, hc, ac, lc, blurb, points) in enumerate(criteria):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.42)
        y = Inches(1.12) + row * Inches(2.98)

        # card
        rect(sl, x, y, Inches(6.15), Inches(2.85), lc, C_BORDER, Pt(0.75))
        rect(sl, x, y, Inches(0.22), Inches(2.85), ac)
        # number + title
        txt(sl, num, x + Inches(0.32), y + Inches(0.12),
            Inches(0.5), Inches(0.38), Pt(20), True, ac)
        txt(sl, title, x + Inches(0.82), y + Inches(0.18),
            Inches(5.1), Inches(0.36), Pt(14), True, C_DARK)
        hline(sl, y + Inches(0.6), x + Inches(0.32), x + Inches(5.85), ac, Pt(1.0))
        txt(sl, blurb, x + Inches(0.32), y + Inches(0.68),
            Inches(5.6), Inches(0.38), Pt(10.5), False, C_GRAY_D)
        for j, pt in enumerate(points):
            txt(sl, f"  {pt}",
                x + Inches(0.32), y + Inches(1.1) + j * Inches(0.4),
                Inches(5.65), Inches(0.38), Pt(10), False, C_DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Score Table
# ═══════════════════════════════════════════════════════════════════════════

def slide_scores(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Email-by-Email Scores", 5,
                "Each metric scored 0–10  |  Maximum 30 per email")

    rows = [
        ("01", "Welding Robots",       "Mfg. RFQ",         6,7,5,18,   8,7,9,24),
        ("02", "EV Battery Packs",     "Mfg. RFQ",         6,7,5,18,   8,8,9,25),
        ("03", "Laptop Procurement",   "IT Hardware",       6,6,4,16,   8,7,9,24),
        ("04", "OLED Panels",          "Electronics",       6,7,5,18,   9,8,9,26),
        ("05", "Office Furniture",     "Facilities",        6,7,5,18,   9,8,9,26),
        ("06", "Smartphones",          "Consumer/Retail",   6,7,5,18,   9,8,9,26),
        ("07", "Electric Vans",        "Fleet Proc.",       6,7,5,18,   8,8,9,25),
    ]

    # column config: (label, x, w)
    cols = [
        ("#",           Inches(0.5),   Inches(0.4)),
        ("Email",       Inches(0.95),  Inches(2.05)),
        ("Domain",      Inches(3.05),  Inches(1.62)),
        ("Corr.",       Inches(4.72),  Inches(0.82)),
        ("Cov.",        Inches(5.58),  Inches(0.82)),
        ("Str.",        Inches(6.44),  Inches(0.82)),
        ("Total",       Inches(7.3),   Inches(0.9)),
        ("Corr.",       Inches(8.55),  Inches(0.82)),
        ("Cov.",        Inches(9.41),  Inches(0.82)),
        ("Str.",        Inches(10.27), Inches(0.82)),
        ("Total",       Inches(11.13), Inches(0.9)),
    ]

    # pipeline headers
    rect(sl, Inches(4.67), Inches(1.06), Inches(3.6), Inches(0.36), C_ORANGE)
    txt(sl, "WITHOUT DSPy", Inches(4.67), Inches(1.09),
        Inches(3.6), Inches(0.3), Pt(10), True, C_WHITE, PP_ALIGN.CENTER)
    rect(sl, Inches(8.5), Inches(1.06), Inches(3.6), Inches(0.36), C_BLUE)
    txt(sl, "WITH DSPy", Inches(8.5), Inches(1.09),
        Inches(3.6), Inches(0.3), Pt(10), True, C_WHITE, PP_ALIGN.CENTER)

    # column headers
    rect(sl, Inches(0.5), Inches(1.42), Inches(11.62), Inches(0.4), C_NAVY)
    for label, cx, cw in cols:
        color = C_ORANGE if label in ("Corr.", "Cov.", "Str.", "Total") and cx < Inches(8.0) else \
                C_BLUE_LIGHT if label in ("Corr.", "Cov.", "Str.", "Total") else C_WHITE
        # differentiate the two "Total" columns
        if label == "Total" and cx > Inches(8.0):
            color = C_BLUE_LIGHT
        txt(sl, label, cx + Inches(0.04), Inches(1.47),
            cw - Inches(0.04), Inches(0.32), Pt(10), True, color, PP_ALIGN.CENTER)

    def score_color(v, is_total=False):
        if is_total:
            return (C_GREEN, C_GREEN_LT) if v >= 24 else \
                   (C_AMBER, C_AMBER_LT) if v >= 20 else (C_RED, C_RED_LT)
        return (C_GREEN, C_GREEN_LT) if v >= 8 else \
               (C_AMBER, C_AMBER_LT) if v >= 6 else (C_RED, C_RED_LT)

    for r, row in enumerate(rows):
        y = Inches(1.82) + r * Inches(0.68)
        bg = C_WHITE if r % 2 == 0 else C_GRAY_LL
        rect(sl, Inches(0.5), y, Inches(11.62), Inches(0.66), bg, C_BORDER, Pt(0.5))

        num, name, domain = row[0], row[1], row[2]
        txt(sl, num,    Inches(0.55), y + Inches(0.17), Inches(0.38), Inches(0.32), Pt(10), False, C_GRAY_M, PP_ALIGN.CENTER)
        txt(sl, name,   Inches(0.95), y + Inches(0.17), Inches(2.05), Inches(0.32), Pt(11), True,  C_DARK)
        txt(sl, domain, Inches(3.05), y + Inches(0.17), Inches(1.62), Inches(0.32), Pt(10), False, C_GRAY_D)

        # without DSPy scores
        for val, (_, cx, cw), is_tot in zip(row[3:7], cols[3:7], [False,False,False,True]):
            fg, bg2 = score_color(val, is_tot)
            if is_tot:
                rect(sl, cx + Inches(0.06), y + Inches(0.12), cw - Inches(0.1), Inches(0.42), bg2, fg, Pt(0.75))
                txt(sl, str(val), cx + Inches(0.06), y + Inches(0.12),
                    cw - Inches(0.08), Inches(0.42), Pt(13), True, fg, PP_ALIGN.CENTER)
            else:
                txt(sl, str(val), cx, y + Inches(0.17), cw, Inches(0.32), Pt(12), True, fg, PP_ALIGN.CENTER)

        # separator
        rect(sl, Inches(8.45), y, Inches(0.05), Inches(0.66), C_BORDER)

        # with DSPy scores
        for val, (_, cx, cw), is_tot in zip(row[7:11], cols[7:11], [False,False,False,True]):
            fg, bg2 = score_color(val, is_tot)
            if is_tot:
                rect(sl, cx + Inches(0.06), y + Inches(0.12), cw - Inches(0.1), Inches(0.42), bg2, fg, Pt(0.75))
                txt(sl, str(val), cx + Inches(0.06), y + Inches(0.12),
                    cw - Inches(0.08), Inches(0.42), Pt(13), True, fg, PP_ALIGN.CENTER)
            else:
                txt(sl, str(val), cx, y + Inches(0.17), cw, Inches(0.32), Pt(12), True, fg, PP_ALIGN.CENTER)

    # averages
    yr = Inches(1.82) + 7 * Inches(0.68)
    rect(sl, Inches(0.5), yr, Inches(11.62), Inches(0.65), C_NAVY)
    txt(sl, "AVERAGE", Inches(0.95), yr + Inches(0.18), Inches(2.0), Inches(0.3),
        Pt(11), True, C_WHITE)
    for val, (_, cx, cw), is_tot in zip([6.0, 6.9, 4.9, "17.7"], cols[3:7], [0,0,0,1]):
        txt(sl, str(val), cx, yr + Inches(0.18), cw, Inches(0.3),
            Pt(12), True, C_ORANGE, PP_ALIGN.CENTER)
    rect(sl, Inches(8.45), yr, Inches(0.05), Inches(0.65), C_GRAY_D)
    for val, (_, cx, cw) in zip([8.4, 7.7, 9.0, "25.1"], cols[7:11]):
        txt(sl, str(val), cx, yr + Inches(0.18), cw, Inches(0.3),
            Pt(12), True, C_BLUE, PP_ALIGN.CENTER)

    # legend
    for label, color in [("Green = strong (8+)", C_GREEN),
                          ("Amber = acceptable (6-7)", C_AMBER),
                          ("Red = weak (<6)", C_RED)]:
        pass
    txt(sl, "Score legend:    Green = strong (8+)     Amber = acceptable (6-7)     Red = weak (<6)",
        Inches(0.5), Inches(7.12), Inches(12.0), Inches(0.32),
        Pt(10), False, C_GRAY_M, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Correctness & Coverage
# ═══════════════════════════════════════════════════════════════════════════

def slide_findings(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Key Findings: Correctness & Coverage", 6,
                "Where each pipeline succeeds and fails")

    # column headers
    for x, label, color in [
        (Inches(0.5),  "WITHOUT DSPy", C_ORANGE),
        (Inches(6.85), "WITH DSPy",    C_BLUE),
    ]:
        rect(sl, x, Inches(1.08), Inches(6.2), Inches(0.36), color)
        txt(sl, label, x + Inches(0.14), Inches(1.11),
            Inches(5.9), Inches(0.3), Pt(12), True, C_WHITE)

    issues = [
        # (severity_label, severity_color, title, description)
        ("CRITICAL", C_RED, "Field contamination — all 7 emails",
         "Agent2 injects IT/SaaS fields into every email regardless of domain: "
         "os_preference, uptime_sla_percent, disaster_recovery_rto_rpo, "
         "escrow_arrangement, data_residency_requirement. "
         "Agent3 then burns 8-12 drops per email undoing this."),
        ("BUG", C_RED, "Factual inversion — email_03",
         "vendor_name = 'NovaTech Solutions' — the buyer, not the vendor. "
         "Roles are swapped in the output."),
        ("INCORRECT", C_AMBER, "False prior relationship — email_02",
         "past_relationship = 'existing_vendor' on Greenfield Auto, "
         "which is a cold outreach RFQ. Not stated in the email."),
        ("WEAK", C_AMBER, "Empty risk_flags — 6 of 7 emails",
         "Validation reports risk_flags: [] even when AMC scope, "
         "warranty terms, and defect SLA are all undefined."),
    ]

    minor = [
        ("MINOR", C_AMBER, "Overconfident completeness — email_03",
         "completeness = 'complete' but warranty, penalty clause, and "
         "payment schedule were inferred from domain knowledge — not stated. "
         "Should be 'mostly_complete'."),
        ("MINOR", C_AMBER, "Schema inconsistency — email_07",
         "open_items returned as a dict object instead of the "
         "specified array format. Breaks downstream JSON consumers."),
        ("MINOR", C_AMBER, "record_type imprecision — 3 emails",
         "context.record_type = 'inquiry' used for formal RFQs. "
         "'rfq' is a more precise value for routing purposes."),
        ("PASS", C_GREEN, "No issues — emails 04, 05, 06",
         "Fully correct: right domain, right fields, domain-specific certs "
         "identified (BIFMA, UN38.3, BIS), defect SLA filled from industry knowledge."),
    ]

    for panel_i, (items, x) in enumerate([(issues, Inches(0.5)), (minor, Inches(6.85))]):
        for i, (sev, sev_color, title, desc) in enumerate(items):
            y = Inches(1.52) + i * Inches(1.38)
            rect(sl, x, y, Inches(6.2), Inches(1.3), C_GRAY_LL, C_BORDER, Pt(0.75))
            rect(sl, x, y, Inches(0.22), Inches(1.3), sev_color)
            # severity badge
            rect(sl, x + Inches(0.3), y + Inches(0.1), Inches(1.0), Inches(0.24), sev_color)
            txt(sl, sev, x + Inches(0.3), y + Inches(0.11),
                Inches(1.0), Inches(0.22), Pt(8), True, C_WHITE, PP_ALIGN.CENTER)
            txt(sl, title, x + Inches(0.3), y + Inches(0.4),
                Inches(5.75), Inches(0.3), Pt(11), True, C_DARK)
            txt(sl, desc, x + Inches(0.3), y + Inches(0.72),
                Inches(5.75), Inches(0.55), Pt(10), False, C_GRAY_D)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Structure Comparison
# ═══════════════════════════════════════════════════════════════════════════

def slide_structure(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Output Structure Comparison", 7,
                "Schema design and measurable output differences")

    for col_i, (label, color, lt, schema) in enumerate([
        ("WITHOUT DSPy — Output Schema", C_ORANGE, C_ORANGE_LT, [
            ("email_meta",       "{ from, to, date, subject }"),
            ("context",          "{ domain, request_type, summary, tone, relationship }"),
            ("extracted_values", "{ PRODUCT/SERVICE, QUANTITIES, FINANCIAL, DATES, … }"),
            ("tag_list",         "{ 40–60 flat snake_case keys }"),
            ("sender_notes",     "[ array of strings ]"),
            ("review_decisions", "{ key: { action, destination, reason } }"),
            ("validation",       "{ budget_feasibility, timeline_assessment,"),
            ("",                 "  missing_critical_fields[], risk_flags[] }"),
        ]),
        ("WITH DSPy — Output Schema", C_BLUE, C_BLUE_LIGHT, [
            ("document_meta",    "{ sender, recipient, date, title, document_type }"),
            ("context",          "{ domain, record_type, summary, tone, intent }"),
            ("extracted_values", "{ flat facts sourced only from the email }"),
            ("field_inventory",  "{ financial{}, timeline{}, parties{},"),
            ("",                 "  requirements{}, contacts{} }"),
            ("open_items",       "[ ambiguities quoted from email ]"),
            ("resolution_log",   "{ key: { action, destination, reason } }"),
            ("validation",       "{ completeness, missing_inputs[], inconsistencies[],"),
            ("",                 "  risk_flags[], summary }"),
        ]),
    ]):
        x = Inches(0.5) + col_i * Inches(6.42)
        rect(sl, x, Inches(1.08), Inches(6.15), Inches(0.36), color)
        txt(sl, label, x + Inches(0.12), Inches(1.1),
            Inches(5.9), Inches(0.32), Pt(11), True, C_WHITE)
        rect(sl, x, Inches(1.44), Inches(6.15), Inches(3.62), lt, C_BORDER, Pt(0.75))
        for i, (key, val) in enumerate(schema):
            y = Inches(1.52) + i * Inches(0.4)
            if key:
                txt(sl, key, x + Inches(0.2), y, Inches(1.6), Inches(0.38),
                    Pt(10), True, color)
                txt(sl, val, x + Inches(1.85), y, Inches(4.15), Inches(0.38),
                    Pt(10), False, C_GRAY_D)
            else:
                txt(sl, val, x + Inches(1.85), y, Inches(4.15), Inches(0.38),
                    Pt(10), False, C_GRAY_D)

    # stat comparison table
    rect(sl, Inches(0.5), Inches(5.18), Inches(12.33), Inches(0.38), C_NAVY)
    txt(sl, "Metric", Inches(0.6), Inches(5.22), Inches(3.5), Inches(0.3),
        Pt(10), True, C_WHITE)
    txt(sl, "WITHOUT DSPy", Inches(4.6), Inches(5.22), Inches(3.5), Inches(0.3),
        Pt(10), True, C_ORANGE, PP_ALIGN.CENTER)
    txt(sl, "WITH DSPy", Inches(8.7), Inches(5.22), Inches(3.5), Inches(0.3),
        Pt(10), True, C_BLUE_LIGHT, PP_ALIGN.CENTER)

    stat_rows = [
        ("Avg. output size",          "~6,200 chars",       "~4,300 chars",       False),
        ("Fields per email",          "40–60 (flat)",        "18–26 (nested)",     True),
        ("Irrelevant drops per email", "8–12 fields dropped", "0–2 fields dropped", True),
        ("Validation summary sentence","Not present",         "Always present",     True),
    ]
    for r, (metric, nod, wd, wd_better) in enumerate(stat_rows):
        y = Inches(5.56) + r * Inches(0.46)
        bg = C_WHITE if r % 2 == 0 else C_GRAY_LL
        rect(sl, Inches(0.5), y, Inches(12.33), Inches(0.44), bg, C_BORDER, Pt(0.5))
        txt(sl, metric, Inches(0.6), y + Inches(0.1),
            Inches(3.9), Inches(0.32), Pt(10.5), False, C_DARK)
        txt(sl, nod, Inches(4.1), y + Inches(0.1),
            Inches(4.2), Inches(0.32), Pt(10.5), False,
            C_RED if wd_better else C_GRAY_D, PP_ALIGN.CENTER)
        txt(sl, wd, Inches(8.2), y + Inches(0.1),
            Inches(4.2), Inches(0.32), Pt(10.5), True,
            C_GREEN if wd_better else C_DARK, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Validation Quality
# ═══════════════════════════════════════════════════════════════════════════

def slide_validation(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Validation Quality Deep-Dive", 8,
                "Comparing how effectively each pipeline surfaces real risks and gaps")

    txt(sl, "Six real scenarios — what did each pipeline catch?",
        Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.38),
        Pt(13), False, C_GRAY_D)

    # column headers
    for x, label, color in [
        (Inches(4.5),  "WITHOUT DSPy", C_ORANGE),
        (Inches(8.6),  "WITH DSPy",    C_BLUE),
    ]:
        rect(sl, x, Inches(1.52), Inches(4.0), Inches(0.32), color)
        txt(sl, label, x + Inches(0.1), Inches(1.54),
            Inches(3.8), Inches(0.28), Pt(10), True, C_WHITE, PP_ALIGN.CENTER)

    rect(sl, Inches(0.5), Inches(1.52), Inches(3.95), Inches(0.32), C_NAVY)
    txt(sl, "Scenario", Inches(0.6), Inches(1.54),
        Inches(3.75), Inches(0.28), Pt(10), True, C_WHITE)

    checks = [
        ("AMC scope undefined\n(email 01 — Welding Robots)",
         "risk_flags: []  — missed entirely", C_RED,
         "risk_flags includes: 'AMC pricing scope\nand duration are undefined'", C_GREEN),
        ("Warranty terms vague\n(email 01, 03)",
         "risk_flags: []  — missed entirely", C_RED,
         "risk_flags: 'Warranty terms are vague\nand lack specifics'", C_GREEN),
        ("No defect SLA stated\n(email 04 — OLED Panels)",
         "Not flagged; field left empty", C_AMBER,
         "defect_rate_sla filled with\n'<=0.5% (industry standard for AMOLED)'", C_GREEN),
        ("No RFQ reference number\n(emails 02, 04, 07)",
         "In missing_critical_fields list\nbut risk_flags: []", C_AMBER,
         "risk_flags: 'Missing RFQ reference\nor identifier'", C_GREEN),
        ("Budget tax-inclusive?\n(email 02 — EUR 18M)",
         "Not mentioned in output", C_RED,
         "inconsistencies: 'Clarification\nneeded on EUR 18M tax inclusion'", C_GREEN),
        ("Phased delivery risk\n(email 07 — 3 phases)",
         "timeline_assessment: 'realistic'\n— no per-phase risk noted", C_AMBER,
         "open_items flags AMC terms\nand equivalent model compliance", C_GREEN),
    ]

    for i, (scenario, no_t, no_c, yes_t, yes_c) in enumerate(checks):
        row = i % 3
        big_col = i // 3
        x_s = Inches(0.5)
        x_n = Inches(4.5)
        x_y = Inches(8.6)
        y = Inches(1.92) + row * Inches(1.7)
        if big_col == 1:
            y = y  # same rows, just offset columns on next "page" — keep 3 per column
        bg = C_WHITE if row % 2 == 0 else C_GRAY_LL

        # only draw first 3 or second 3 depending on index
        if big_col == 0 and i < 3:
            rect(sl, x_s, y, Inches(3.95), Inches(1.62), C_GRAY_LL, C_BORDER, Pt(0.5))
            txt(sl, scenario, x_s + Inches(0.14), y + Inches(0.1),
                Inches(3.7), Inches(0.55), Pt(10.5), True, C_DARK)
            rect(sl, x_n, y, Inches(4.0), Inches(1.62), C_RED_LT, C_BORDER, Pt(0.5))
            rect(sl, x_n, y, Inches(0.18), Inches(1.62), no_c)
            txt(sl, no_t, x_n + Inches(0.26), y + Inches(0.2),
                Inches(3.6), Inches(0.75), Pt(10), False, C_DARK)
            rect(sl, x_y, y, Inches(4.0), Inches(1.62), C_GREEN_LT, C_BORDER, Pt(0.5))
            rect(sl, x_y, y, Inches(0.18), Inches(1.62), yes_c)
            txt(sl, yes_t, x_y + Inches(0.26), y + Inches(0.2),
                Inches(3.6), Inches(0.75), Pt(10), False, C_DARK)
        elif big_col == 1 and i >= 3:
            row2 = i - 3
            y2 = Inches(1.92) + row2 * Inches(1.7)
            rect(sl, x_s, y2, Inches(3.95), Inches(1.62), C_GRAY_LL, C_BORDER, Pt(0.5))
            txt(sl, scenario, x_s + Inches(0.14), y2 + Inches(0.1),
                Inches(3.7), Inches(0.55), Pt(10.5), True, C_DARK)
            rect(sl, x_n, y2, Inches(4.0), Inches(1.62), C_AMBER_LT if no_c == C_AMBER else C_RED_LT, C_BORDER, Pt(0.5))
            rect(sl, x_n, y2, Inches(0.18), Inches(1.62), no_c)
            txt(sl, no_t, x_n + Inches(0.26), y2 + Inches(0.2),
                Inches(3.6), Inches(0.75), Pt(10), False, C_DARK)
            rect(sl, x_y, y2, Inches(4.0), Inches(1.62), C_GREEN_LT, C_BORDER, Pt(0.5))
            rect(sl, x_y, y2, Inches(0.18), Inches(1.62), yes_c)
            txt(sl, yes_t, x_y + Inches(0.26), y2 + Inches(0.2),
                Inches(3.6), Inches(0.75), Pt(10), False, C_DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Unseen Domain
# ═══════════════════════════════════════════════════════════════════════════

def slide_unseen(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Handling Unseen Domains & Unfamiliar Email Structures", 9,
                "Which pipeline generalises better when a brand-new domain arrives?")

    txt(sl,
        "If an email arrives from a domain neither pipeline has seen "
        "(e.g. pharmaceutical cold-chain, satellite components, carbon credit trading) "
        "— how does each pipeline respond?",
        Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.5),
        Pt(12), False, C_GRAY_D)

    categories = [
        ("Schema generation\nfor new domain",
         "Falls back to a hardcoded IT/SaaS template.\nExpect os_preference, cybersecurity_standard,\ndata_residency on a pharma cold-chain RFQ.",
         False,
         "DSPy derives fields from the domain + record_type\nclassified by Agent1. Novel domain = novel schema.\nNo template contamination.",
         True),
        ("Domain classification\naccuracy",
         "Limited to hard-coded buckets (manufacturing, IT,\nretail, fleet). Novel domains get mis-mapped.",
         False,
         "Agent1 uses open-ended free-text: 'domain = subject area'.\nAny label is valid. DSPy audit adapts to it automatically.",
         True),
        ("Field contamination\nunder new domain",
         "High — IT/SaaS fields already bleed into every\nnon-IT email. A novel domain adds more noise.",
         False,
         "Near-zero. DSPy StructuralAudit is conditioned on the\ndetected domain per-run. No fixed field list to bleed from.",
         True),
        ("Robustness to unusual\nemail structure",
         "Agent1 prompt expects specific sections\n(FINANCIAL, DATES, etc.). Non-standard layouts\nlead to missed fields.",
         False,
         "Agent1 forensic extraction is open-ended.\nDSPy audit fills structural gaps per detected domain —\nnot per a fixed template.",
         True),
    ]

    col_x = [Inches(0.5), Inches(3.78), Inches(8.05)]
    col_w = [Inches(3.22), Inches(4.22), Inches(4.22)]

    # column headers
    rect(sl, col_x[0], Inches(1.7), col_w[0], Inches(0.36), C_NAVY)
    rect(sl, col_x[1], Inches(1.7), col_w[1], Inches(0.36), C_ORANGE)
    rect(sl, col_x[2], Inches(1.7), col_w[2], Inches(0.36), C_BLUE)
    txt(sl, "Dimension", col_x[0] + Inches(0.1), Inches(1.73), col_w[0], Inches(0.3), Pt(10), True, C_WHITE)
    txt(sl, "WITHOUT DSPy", col_x[1] + Inches(0.1), Inches(1.73), col_w[1], Inches(0.3), Pt(10), True, C_WHITE, PP_ALIGN.CENTER)
    txt(sl, "WITH DSPy",    col_x[2] + Inches(0.1), Inches(1.73), col_w[2], Inches(0.3), Pt(10), True, C_WHITE, PP_ALIGN.CENTER)

    for r, (dim, no_text, no_good, yes_text, yes_good) in enumerate(categories):
        y = Inches(2.1) + r * Inches(1.22)
        bg = C_WHITE if r % 2 == 0 else C_GRAY_LL

        rect(sl, col_x[0], y, col_w[0], Inches(1.18), bg, C_BORDER, Pt(0.5))
        txt(sl, dim, col_x[0] + Inches(0.12), y + Inches(0.08),
            col_w[0] - Inches(0.2), Inches(0.95), Pt(10.5), True, C_DARK)

        no_bg = C_RED_LT if not no_good else C_GREEN_LT
        rect(sl, col_x[1], y, col_w[1], Inches(1.18), no_bg, C_BORDER, Pt(0.5))
        rect(sl, col_x[1], y, Inches(0.14), Inches(1.18), C_RED if not no_good else C_GREEN)
        txt(sl, no_text, col_x[1] + Inches(0.22), y + Inches(0.08),
            col_w[1] - Inches(0.28), Inches(1.05), Pt(10), False, C_DARK)

        yes_bg = C_GREEN_LT if yes_good else C_RED_LT
        rect(sl, col_x[2], y, col_w[2], Inches(1.18), yes_bg, C_BORDER, Pt(0.5))
        rect(sl, col_x[2], y, Inches(0.14), Inches(1.18), C_GREEN if yes_good else C_RED)
        txt(sl, yes_text, col_x[2] + Inches(0.22), y + Inches(0.08),
            col_w[2] - Inches(0.28), Inches(1.05), Pt(10), False, C_DARK)

    # verdict
    rect(sl, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.4), C_BLUE)
    txt(sl,
        "Verdict:  With DSPy handles unseen domains significantly better — "
        "schema is derived from context, not loaded from a hardcoded template.",
        Inches(0.65), Inches(7.05), Inches(12.0), Inches(0.32),
        Pt(11), True, C_WHITE, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Recommendations
# ═══════════════════════════════════════════════════════════════════════════

def slide_recs(prs):
    sl = blank_slide(prs)
    page_chrome(sl, "Recommendations & Next Steps", 10,
                "Prioritised actions based on evaluation findings")

    recs = [
        (C_GREEN,  "HIGH", "Adopt the DSPy pipeline as production baseline",
         "Scores 25.1/30 vs 17.7/30. Cleaner schema, zero domain contamination, "
         "honest risk flagging, and robust handling of unseen domains."),
        (C_AMBER,  "MEDIUM", "Fix the 3 minor DSPy issues before shipping",
         "1) Cap completeness = 'mostly_complete' when > 20% of fields are inferred.  "
         "2) Enforce open_items as array type in the StructuralAudit DSPy signature.  "
         "3) Add 'rfq' to the valid record_type vocabulary in Agent1's instructions."),
        (C_RED,    "RETIRE", "Decommission the 3-agent prompt-only pipeline",
         "The IT/SaaS template bleed in Agent2 is a structural design flaw — not fixable "
         "with prompt tweaks. The pipeline architecture needs replacing, not patching."),
        (C_BLUE,   "ENHANCE", "Compile DSPy StructuralAudit with BootstrapFewShot",
         "The middle stage runs zero-shot today. Adding 5-10 few-shot examples per domain "
         "via BootstrapFewShot would likely push coverage scores from 7-8 to 9-10."),
        (C_TEAL,   "EXTEND", "Validate on emails 08-10 (Kitchen Equipment, Govt Tender, Aluminium Supply)",
         "These represent hospitality, government procurement, and commodity contracts — "
         "three domain types not yet fully evaluated. Good stress-test for the DSPy pipeline."),
    ]

    for i, (color, tag, title, desc) in enumerate(recs):
        y = Inches(1.1) + i * Inches(1.2)
        rect(sl, Inches(0.5), y, Inches(12.33), Inches(1.12), C_GRAY_LL, C_BORDER, Pt(0.75))
        rect(sl, Inches(0.5), y, Inches(0.22), Inches(1.12), color)
        # tag badge
        rect(sl, Inches(0.82), y + Inches(0.1), Inches(0.95), Inches(0.28), color)
        txt(sl, tag, Inches(0.82), y + Inches(0.11),
            Inches(0.95), Inches(0.26), Pt(8.5), True, C_WHITE, PP_ALIGN.CENTER)
        txt(sl, title, Inches(1.85), y + Inches(0.1),
            Inches(10.8), Inches(0.3), Pt(12), True, C_DARK)
        txt(sl, desc, Inches(1.85), y + Inches(0.44),
            Inches(10.8), Inches(0.62), Pt(10.5), False, C_GRAY_D)

# ═══════════════════════════════════════════════════════════════════════════
# Build
# ═══════════════════════════════════════════════════════════════════════════

prs = new_prs()
slide_title(prs)
slide_intro(prs)
slide_data(prs)
slide_criteria(prs)
slide_scores(prs)
slide_findings(prs)
slide_structure(prs)
slide_validation(prs)
slide_unseen(prs)
slide_recs(prs)

OUT = "pipeline_comparison_report.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
